from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Colour palette ──────────────────────────────────────────────────────────
TEAL        = RGBColor(0x0f, 0x76, 0x6e)
TEAL_LIGHT  = RGBColor(0x14, 0xb8, 0xa6)
TEAL_DARK   = RGBColor(0x05, 0x0f, 0x0e)
DARK_BG     = RGBColor(0x0a, 0x0e, 0x0d)
NEAR_BLACK  = RGBColor(0x06, 0x09, 0x09)
WHITE       = RGBColor(0xff, 0xff, 0xff)
AMBER       = RGBColor(0xf5, 0x9e, 0x0b)
RED         = RGBColor(0xef, 0x44, 0x44)
PURPLE      = RGBColor(0xa7, 0x8b, 0xfa)
GRAY_DIM    = RGBColor(0x64, 0x74, 0x8b)
GRAY_MID    = RGBColor(0x33, 0x41, 0x55)
CARD_BG     = RGBColor(0x12, 0x1e, 0x1c)
CARD_BORDER = RGBColor(0x1f, 0x3b, 0x38)

W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

BLANK = prs.slide_layouts[6]   # completely blank

# ── Helpers ──────────────────────────────────────────────────────────────────

def add_slide():
    return prs.slides.add_slide(BLANK)

def bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def rect(slide, x, y, w, h, fill_color=None, line_color=None, line_width=Pt(0)):
    from pptx.util import Emu
    shape = slide.shapes.add_shape(1, x, y, w, h)   # MSO_SHAPE_TYPE.RECTANGLE = 1
    shape.line.width = line_width
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = line_width if line_width else Pt(0.75)
    else:
        shape.line.fill.background()
    return shape

def txbox(slide, text, x, y, w, h,
          size=Pt(14), bold=False, color=WHITE,
          align=PP_ALIGN.LEFT, wrap=True, italic=False):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    p  = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = size
    run.font.bold  = bold
    run.font.color.rgb = color
    run.font.italic = italic
    return tb

def eyebrow(slide, text, x=Inches(0.7), y=Inches(0.45)):
    txbox(slide, text, x, y, Inches(11), Inches(0.35),
          size=Pt(10), bold=True, color=TEAL_LIGHT)

def slide_title(slide, text, x=Inches(0.7), y=Inches(0.85), w=Inches(11.9), h=Inches(1.1)):
    txbox(slide, text, x, y, w, h, size=Pt(30), bold=True, color=WHITE)

def bullet(slide, icon, title, body, x, y, w=Inches(11.9), card_h=Inches(1.05)):
    r = rect(slide, x, y, w, card_h, fill_color=CARD_BG, line_color=CARD_BORDER, line_width=Pt(0.5))
    # icon
    txbox(slide, icon, x + Inches(0.18), y + Inches(0.18),
          Inches(0.55), card_h - Inches(0.18), size=Pt(20))
    # title
    txbox(slide, title, x + Inches(0.85), y + Inches(0.1),
          w - Inches(1.05), Inches(0.38), size=Pt(13), bold=True, color=WHITE)
    # body
    txbox(slide, body, x + Inches(0.85), y + Inches(0.45),
          w - Inches(1.05), Inches(0.55), size=Pt(11), color=GRAY_DIM)

def stat_card(slide, value, label, x, y, w=Inches(3.7), h=Inches(2.0), accent=TEAL_LIGHT):
    rect(slide, x, y, w, h, fill_color=CARD_BG, line_color=CARD_BORDER, line_width=Pt(0.5))
    # accent top bar
    rect(slide, x, y, w, Pt(3), fill_color=accent)
    txbox(slide, value, x + Inches(0.2), y + Inches(0.25),
          w - Inches(0.3), Inches(0.8), size=Pt(36), bold=True, color=accent)
    txbox(slide, label, x + Inches(0.2), y + Inches(1.0),
          w - Inches(0.3), Inches(0.9), size=Pt(12), color=GRAY_DIM, wrap=True)

def section_slide(num, label, title_text):
    sl = add_slide()
    bg(sl, NEAR_BLACK)
    # big faded number
    txbox(sl, num, Inches(0), Inches(1.2), W, Inches(5),
          size=Pt(200), bold=True, color=RGBColor(0x0d,0x4f,0x4a),
          align=PP_ALIGN.CENTER)
    # label
    txbox(sl, label, Inches(1), Inches(2.4), Inches(11.3), Inches(0.5),
          size=Pt(12), bold=True, color=TEAL_LIGHT, align=PP_ALIGN.CENTER)
    # title
    txbox(sl, title_text, Inches(1), Inches(3.0), Inches(11.3), Inches(2.0),
          size=Pt(54), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    return sl

def highlight_box(slide, title, body, x, y, w, h, accent=TEAL_LIGHT, border=TEAL):
    rect(slide, x, y, w, h, fill_color=CARD_BG, line_color=border, line_width=Pt(0.75))
    txbox(slide, title, x + Inches(0.2), y + Inches(0.15),
          w - Inches(0.3), Inches(0.35), size=Pt(12), bold=True, color=accent)
    txbox(slide, body, x + Inches(0.2), y + Inches(0.5),
          w - Inches(0.3), h - Inches(0.6), size=Pt(11), color=GRAY_DIM, wrap=True)


# ════════════════════════════════════════════════════════════════════════════
#  SLIDE 1 — COVER
# ════════════════════════════════════════════════════════════════════════════
sl = add_slide()
bg(sl, NEAR_BLACK)
rect(sl, Inches(0), Inches(0), W, H, fill_color=RGBColor(0x05, 0x12, 0x10))

# Gradient-like teal glow (faked with a large transparent-ish rect)
glow = rect(sl, Inches(0), Inches(0), Inches(7), Inches(5),
            fill_color=RGBColor(0x0d, 0x4f, 0x4a))
glow.fill.fore_color.theme_color  # access to confirm
# Badge
rect(sl, Inches(4.4), Inches(1.35), Inches(4.5), Inches(0.38),
     fill_color=RGBColor(0x0a,0x2e,0x2b), line_color=TEAL, line_width=Pt(0.75))
txbox(sl, "MENTAL HEALTH TECH  ·  NEPAL  ·  2026",
      Inches(4.5), Inches(1.38), Inches(4.3), Inches(0.32),
      size=Pt(9), bold=True, color=TEAL_LIGHT, align=PP_ALIGN.CENTER)

txbox(sl, "AAMA", Inches(1.5), Inches(1.85), Inches(10), Inches(2.0),
      size=Pt(110), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
txbox(sl, "आमा सखी", Inches(1.5), Inches(3.5), Inches(10), Inches(0.8),
      size=Pt(32), bold=True, color=TEAL_LIGHT, align=PP_ALIGN.CENTER)
txbox(sl,
      "An offline-first mobile platform empowering Nepal's 50,000+ Female Community\n"
      "Health Volunteers to identify and triage mental health crises — before they become emergencies.",
      Inches(2.0), Inches(4.35), Inches(9.3), Inches(1.1),
      size=Pt(13), color=GRAY_DIM, align=PP_ALIGN.CENTER, wrap=True)

tags = ["React Native · Expo", "EPDS · PHQ-A Screening", "Offline-First SQLite", "Bilingual Nepali/English"]
for i, t in enumerate(tags):
    tx = Inches(0.7) + i * Inches(3.1)
    rect(sl, tx, Inches(5.6), Inches(2.9), Inches(0.38),
         fill_color=RGBColor(0x12,0x1e,0x1c), line_color=GRAY_MID, line_width=Pt(0.5))
    txbox(sl, t, tx + Inches(0.1), Inches(5.62), Inches(2.7), Inches(0.34),
          size=Pt(10), color=GRAY_DIM, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════════════════
#  SLIDE 2 — SECTION: PROBLEM
# ════════════════════════════════════════════════════════════════════════════
section_slide("01", "THE PROBLEM", "A Crisis Hidden\nin Plain Sight")


# ════════════════════════════════════════════════════════════════════════════
#  SLIDE 3 — PROBLEM STATS
# ════════════════════════════════════════════════════════════════════════════
sl = add_slide()
bg(sl, DARK_BG)
eyebrow(sl, "PROBLEM  ·  NEPAL'S MENTAL HEALTH GAP")
slide_title(sl, "Vulnerable populations go undiagnosed —\nby design of the existing system")

stat_card(sl, "50,000+",
          "FCHVs conducting routine home visits across all 77 districts with no mental health screening tools",
          Inches(0.7), Inches(2.1), accent=TEAL_LIGHT)
stat_card(sl, "1 in 5",
          "Postnatal women in rural Nepal experience postpartum depression — majority undetected & untreated",
          Inches(4.55), Inches(2.1), accent=AMBER)
stat_card(sl, "< 1%",
          "Of Nepal's health budget allocated to mental health — despite depression being the leading cause of disability",
          Inches(8.4), Inches(2.1), accent=RED)

# quote bar
rect(sl, Inches(0.7), Inches(4.35), Pt(3), Inches(1.1), fill_color=TEAL_LIGHT)
rect(sl, Inches(0.7), Inches(4.35), Inches(11.9), Inches(1.1),
     fill_color=RGBColor(0x0a,0x1e,0x1b), line_color=TEAL, line_width=Pt(0.5))
txbox(sl,
      '"FCHVs are Nepal\'s most trusted healthcare touchpoint — yet they have no structured way\n'
      'to identify or escalate mental health distress."',
      Inches(1.0), Inches(4.45), Inches(11.3), Inches(0.75),
      size=Pt(12), italic=True, color=RGBColor(0xcc,0xff,0xf1))


# ════════════════════════════════════════════════════════════════════════════
#  SLIDE 4 — ROOT CAUSES
# ════════════════════════════════════════════════════════════════════════════
sl = add_slide()
bg(sl, DARK_BG)
eyebrow(sl, "PROBLEM  ·  ROOT CAUSES")
slide_title(sl, "Three systemic failures compounding the crisis")

bullet(sl, "🚫", "No Screening Tools in the Field",
       "FCHVs conduct physical health checks but have never had validated, translated mental health screening instruments (EPDS, PHQ-A) or training to use them.",
       Inches(0.7), Inches(2.0))
bullet(sl, "📡", "Connectivity Deserts",
       "Large swaths of Nepal's rural terrain have no mobile data or internet. Existing digital health tools fail entirely in these environments — data is never recorded.",
       Inches(0.7), Inches(3.2))
bullet(sl, "⚡", "No Escalation Pathway",
       "Even when a volunteer suspects a crisis, there is no formal fast pathway to connect a patient to a psychiatrist. High-risk cases go unescalated — silently.",
       Inches(0.7), Inches(4.4))


# ════════════════════════════════════════════════════════════════════════════
#  SLIDE 5 — SECTION: SOLUTION
# ════════════════════════════════════════════════════════════════════════════
section_slide("02", "THE SOLUTION", "A Companion in\nEvery Home Visit")


# ════════════════════════════════════════════════════════════════════════════
#  SLIDE 6 — SOLUTION DETAIL
# ════════════════════════════════════════════════════════════════════════════
sl = add_slide()
bg(sl, DARK_BG)
eyebrow(sl, "SOLUTION  ·  AAMA सखी")
slide_title(sl, "Clinical-grade screening in the hands\nof community volunteers")

# Left column — what it does
txbox(sl, "WHAT IT DOES", Inches(0.7), Inches(2.0), Inches(5.8), Inches(0.3),
      size=Pt(9), bold=True, color=GRAY_DIM)
bullet(sl, "🩺", "Guided 4-Step Visit Flow",
       "Patient info → physical checks → EPDS/PHQ-A screening → instant risk score & action plan",
       Inches(0.7), Inches(2.35), w=Inches(5.8), card_h=Inches(1.0))
bullet(sl, "🔔", "Auto-Triage & Psychiatrist Routing",
       "High/critical cases auto-matched to nearest psychiatrist by district. SOS overlay for self-harm.",
       Inches(0.7), Inches(3.45), w=Inches(5.8), card_h=Inches(1.0))
bullet(sl, "📴", "100% Offline Operation",
       "All data stored locally in SQLite. Works in zero-connectivity villages. Syncs when online.",
       Inches(0.7), Inches(4.55), w=Inches(5.8), card_h=Inches(1.0))

# Right column — key innovations
txbox(sl, "KEY INNOVATION", Inches(6.8), Inches(2.0), Inches(5.8), Inches(0.3),
      size=Pt(9), bold=True, color=GRAY_DIM)
highlight_box(sl,
    "Transparent Clinical Approach",
    "Volunteers explicitly disclose the mental health check — using a warm, reassuring tone. Questions are shown as validated (no masking), building genuine patient trust.",
    Inches(6.8), Inches(2.35), Inches(5.8), Inches(1.5))
highlight_box(sl,
    "Bilingual by Default",
    "Nepali-first UI with English toggle. Every string, question, and instruction is fully localized — ensuring accessibility for both volunteers and patients.",
    Inches(6.8), Inches(4.0), Inches(5.8), Inches(1.55))


# ════════════════════════════════════════════════════════════════════════════
#  SLIDE 7 — SECTION: TECHNICAL
# ════════════════════════════════════════════════════════════════════════════
section_slide("03", "TECHNICAL COMPLEXITY", "Built to Work\nWhere Nothing Else Does")


# ════════════════════════════════════════════════════════════════════════════
#  SLIDE 8 — TECH ARCHITECTURE
# ════════════════════════════════════════════════════════════════════════════
sl = add_slide()
bg(sl, DARK_BG)
eyebrow(sl, "TECHNICAL  ·  ARCHITECTURE")
slide_title(sl, "Data flow: home visit → risk score → psychiatrist alert")

# Flow steps
steps = [
    ("1", "Patient Info", "Name, age,\ntype, location"),
    ("2", "Physical Checks", "Vitals &\nhealth flags"),
    ("3", "Screening Qs", "EPDS or\nPHQ-A"),
    ("4", "Risk Engine", "Score →\nlevel"),
    ("5", "Alert + Route", "Nearest\npsychiatrist"),
]
box_w = Inches(2.1)
box_h = Inches(1.3)
gap   = Inches(0.22)
start_x = Inches(0.5)
y_flow = Inches(2.1)
for i, (num, title, sub) in enumerate(steps):
    x = start_x + i * (box_w + gap)
    is_last = (i == len(steps) - 1)
    bc = RGBColor(0x18,0x3d,0x39) if not is_last else RGBColor(0x3b,0x0c,0x0c)
    lc = TEAL if not is_last else RED
    rect(sl, x, y_flow, box_w, box_h, fill_color=bc, line_color=lc, line_width=Pt(0.75))
    txbox(sl, f"STEP {num}", x + Inches(0.1), y_flow + Inches(0.08),
          box_w - Inches(0.15), Inches(0.22),
          size=Pt(8), bold=True, color=TEAL_LIGHT if not is_last else RED)
    txbox(sl, title, x + Inches(0.1), y_flow + Inches(0.3),
          box_w - Inches(0.15), Inches(0.4),
          size=Pt(13), bold=True)
    txbox(sl, sub, x + Inches(0.1), y_flow + Inches(0.7),
          box_w - Inches(0.15), Inches(0.5),
          size=Pt(10), color=GRAY_DIM)
    if i < len(steps) - 1:
        txbox(sl, "→", x + box_w + Inches(0.02), y_flow + Inches(0.35),
              Inches(0.18), Inches(0.6), size=Pt(18), color=TEAL_LIGHT, align=PP_ALIGN.CENTER)

# Tech cards (2x2)
tech = [
    ("FRONTEND", "React Native  ·  Expo SDK  ·  React Navigation  ·  AsyncStorage"),
    ("STORAGE",  "expo-sqlite  ·  Local SQLite DB  ·  JSON columns  ·  Offline-first"),
    ("CLINICAL LOGIC", "riskEngine.js  ·  EPDS (0–30)  ·  PHQ-A (0–27)  ·  Self-harm override"),
    ("CROSS-PLATFORM", "iOS  ·  Android  ·  Web (Metro)  ·  Platform shims (.web.js)"),
]
card_w = Inches(5.9)
card_h = Inches(1.15)
for i, (title, pills) in enumerate(tech):
    col = i % 2
    row = i // 2
    cx = Inches(0.7) + col * (card_w + Inches(0.6))
    cy = Inches(3.65) + row * (card_h + Inches(0.14))
    rect(sl, cx, cy, card_w, card_h, fill_color=CARD_BG, line_color=CARD_BORDER, line_width=Pt(0.5))
    txbox(sl, title, cx + Inches(0.18), cy + Inches(0.1),
          card_w - Inches(0.3), Inches(0.3),
          size=Pt(9), bold=True, color=TEAL_LIGHT)
    txbox(sl, pills, cx + Inches(0.18), cy + Inches(0.42),
          card_w - Inches(0.3), Inches(0.6),
          size=Pt(11), color=GRAY_DIM)


# ════════════════════════════════════════════════════════════════════════════
#  SLIDE 9 — RISK ENGINE
# ════════════════════════════════════════════════════════════════════════════
sl = add_slide()
bg(sl, DARK_BG)
eyebrow(sl, "TECHNICAL  ·  CLINICAL INTELLIGENCE")
slide_title(sl, "Validated screening with automatic risk stratification")

# Risk table header
table_x = Inches(0.7)
table_y = Inches(2.0)
col_widths = [Inches(2.3), Inches(1.5), Inches(1.6), Inches(1.5), Inches(1.9)]
headers = ["Scale", "Low", "Moderate", "High", "Critical"]
hx = table_x
for i, (h, cw) in enumerate(zip(headers, col_widths)):
    rect(sl, hx, table_y, cw, Inches(0.38),
         fill_color=RGBColor(0x0d,0x1e,0x1c), line_color=CARD_BORDER, line_width=Pt(0.5))
    txbox(sl, h, hx + Inches(0.1), table_y + Inches(0.06),
          cw - Inches(0.15), Inches(0.28), size=Pt(10), bold=True, color=GRAY_DIM)
    hx += cw

rows = [
    ("EPDS  (Postnatal)",  "0–9",  "10–12", "13+",  "Q10 ★"),
    ("PHQ-A  Ages 12–14", "—",    "—",     "≥ 13", "Q9  ★"),
    ("PHQ-A  Ages 15–19", "—",    "—",     "≥ 11", "Q9  ★"),
]
risk_colors = [
    RGBColor(0x4a,0xde,0x80),
    RGBColor(0xfb,0xbf,0x24),
    RGBColor(0xfb,0x92,0x3c),
    RGBColor(0xf8,0x71,0x71),
]
for ri, row in enumerate(rows):
    ry = table_y + Inches(0.38) + ri * Inches(0.52)
    rx = table_x
    for ci, (cell, cw) in enumerate(zip(row, col_widths)):
        rect(sl, rx, ry, cw, Inches(0.52),
             fill_color=RGBColor(0x0f,0x1a,0x18), line_color=CARD_BORDER, line_width=Pt(0.35))
        color = WHITE if ci == 0 else (GRAY_DIM if cell == "—" else risk_colors[ci - 1])
        txbox(sl, cell, rx + Inches(0.1), ry + Inches(0.1),
              cw - Inches(0.15), Inches(0.34),
              size=Pt(11), bold=(ci > 0 and cell != "—"), color=color)
        rx += cw

txbox(sl, "★  Self-harm response triggers CRITICAL override regardless of total score",
      table_x, table_y + Inches(0.38) + 3 * Inches(0.52) + Inches(0.1),
      Inches(9), Inches(0.3), size=Pt(10), color=RGBColor(0xf8,0x71,0x71))

# Right column features
highlight_box(sl,
    "🚨  SOSOverlay — Non-Dismissable Alert",
    "Fullscreen critical alert on self-harm detection. Cannot be accidentally closed — requires intentional acknowledgment before the volunteer can continue.",
    Inches(7.8), Inches(2.0), Inches(4.8), Inches(1.7),
    accent=RED, border=RED)
highlight_box(sl,
    "📍  District-Based Psychiatrist Routing",
    "alertService.js matches high/critical visits against a psychiatrist roster by district. Auto-assigns nearest available specialist with upsert logic.",
    Inches(7.8), Inches(3.85), Inches(4.8), Inches(1.6))


# ════════════════════════════════════════════════════════════════════════════
#  SLIDE 10 — SECTION: IMPACT
# ════════════════════════════════════════════════════════════════════════════
section_slide("04", "BUSINESS IMPACT & FEASIBILITY", "Scalable Without\na Backend")


# ════════════════════════════════════════════════════════════════════════════
#  SLIDE 11 — IMPACT & FEASIBILITY
# ════════════════════════════════════════════════════════════════════════════
sl = add_slide()
bg(sl, DARK_BG)
eyebrow(sl, "IMPACT  ·  SCALE & FEASIBILITY")
slide_title(sl, "Built to deploy across Nepal's entire FCHV network")

impact_items = [
    ("50K+",  "FCHVs who can deploy this app immediately",                 TEAL_LIGHT),
    ("77",    "Districts covered via psychiatrist roster & district routing", TEAL_LIGHT),
    ("0",     "Internet required during a patient visit",                  AMBER),
    ("$0",    "Infrastructure cost per deployment — fully device-local",   AMBER),
]
iw = Inches(2.75)
ih = Inches(1.9)
for i, (val, label, color) in enumerate(impact_items):
    ix = Inches(0.7) + i * (iw + Inches(0.27))
    rect(sl, ix, Inches(2.0), iw, ih, fill_color=CARD_BG, line_color=CARD_BORDER, line_width=Pt(0.5))
    rect(sl, ix, Inches(2.0), iw, Pt(3), fill_color=color)
    txbox(sl, val, ix + Inches(0.15), Inches(2.2), iw - Inches(0.2), Inches(0.7),
          size=Pt(40), bold=True, color=color)
    txbox(sl, label, ix + Inches(0.15), Inches(2.95), iw - Inches(0.2), Inches(0.8),
          size=Pt(11), color=GRAY_DIM, wrap=True)

highlight_box(sl,
    "Minimal Deployment Barrier",
    "No server, no cloud subscription, no ongoing cost. Distribute as an APK to any Android device. FCHVs already carry smartphones as part of Nepal's DHIS2 rollout.",
    Inches(0.7), Inches(4.1), Inches(5.9), Inches(1.7))
highlight_box(sl,
    "Integration Pathway",
    "Sync flag on every visit record. Designed to plug into Nepal's existing DHIS2 health data infrastructure when connectivity is available — no re-architecture needed.",
    Inches(6.8), Inches(4.1), Inches(5.8), Inches(1.7),
    accent=AMBER, border=RGBColor(0xd9,0x77,0x06))


# ════════════════════════════════════════════════════════════════════════════
#  SLIDE 12 — CHALLENGES
# ════════════════════════════════════════════════════════════════════════════
sl = add_slide()
bg(sl, DARK_BG)
eyebrow(sl, "CHALLENGES  ·  ENGINEERING & ETHICAL")
slide_title(sl, "Hard problems we solved — and some still ahead")

challenges = [
    ("c1", RED,         "Connectivity-Zero Environments",
     "Most rural visit sites have no mobile data. Standard data-centric architectures fail entirely.",
     "Solved: expo-sqlite + .web.js platform shim for full offline-first operation with zero dependency on network state."),
    ("c2", AMBER,       "Cross-Platform Parity",
     "expo-sqlite has no web module — throws at runtime on web build. Metro 500 errors hide the real cause.",
     "Solved: Metro platform-suffix resolution — index.web.js substitutes native DB with an in-memory API-compatible layer."),
    ("c3", TEAL_LIGHT,  "Clinical Accuracy vs. Usability",
     "Screening tools must remain validated — masking or paraphrasing questions invalidates them clinically.",
     "Approach: Transparent disclosure model — volunteers explain purpose first, then administer verbatim Nepali questions."),
    ("c4", PURPLE,      "Self-Harm Response Protocol",
     "Digital tools must not create a false sense of safety. A flagged Q10/Q9 needs immediate, unavoidable human escalation.",
     "Solved: SOSOverlay is fullscreen & non-dismissable — forces acknowledgment before the volunteer can continue."),
]
cw = Inches(5.75)
ch = Inches(2.15)
for i, (_, accent, ctitle, problem, solution) in enumerate(challenges):
    col = i % 2
    row = i // 2
    cx = Inches(0.7) + col * (cw + Inches(0.6))
    cy = Inches(2.05) + row * (ch + Inches(0.18))
    rect(sl, cx, cy, cw, ch, fill_color=CARD_BG, line_color=CARD_BORDER, line_width=Pt(0.5))
    rect(sl, cx, cy, Pt(3), ch, fill_color=accent)
    txbox(sl, ctitle, cx + Inches(0.18), cy + Inches(0.12),
          cw - Inches(0.28), Inches(0.35), size=Pt(13), bold=True)
    txbox(sl, problem, cx + Inches(0.18), cy + Inches(0.5),
          cw - Inches(0.28), Inches(0.55), size=Pt(11), color=GRAY_DIM, wrap=True)
    txbox(sl, solution, cx + Inches(0.18), cy + Inches(1.1),
          cw - Inches(0.28), Inches(0.88), size=Pt(11), color=WHITE, wrap=True)


# ════════════════════════════════════════════════════════════════════════════
#  SLIDE 13 — ROADMAP
# ════════════════════════════════════════════════════════════════════════════
sl = add_slide()
bg(sl, DARK_BG)
eyebrow(sl, "IMPACT  ·  ROADMAP")
slide_title(sl, "From MVP to national deployment")

phases = [
    (TEAL_LIGHT, "MVP  ·  Now",
     "Core Screening Platform",
     "Full EPDS + PHQ-A screening  ·  Offline SQLite storage  ·  Risk engine  ·  SOS overlay  ·  District psychiatrist routing  ·  Bilingual UI"),
    (AMBER, "Phase 2  ·  Q3 2026",
     "DHIS2 Sync & Ministry Partnership",
     "Background sync to Nepal DHIS2  ·  Pilot with 100 FCHVs in Kaski & Lalitpur  ·  Ministry of Health data sharing MOU"),
    (PURPLE, "Phase 3  ·  2027",
     "National Scale & Additional Instruments",
     "All 77 districts  ·  GAD-7 for general anxiety  ·  ASSIST for substance use  ·  FCHV training module  ·  Aggregate policy dashboards"),
]
phase_h = Inches(1.4)
gap_p   = Inches(0.18)
for i, (color, phase_label, ptitle, items) in enumerate(phases):
    py = Inches(2.05) + i * (phase_h + gap_p)
    # phase label
    txbox(sl, phase_label, Inches(0.4), py + Inches(0.15),
          Inches(1.7), Inches(0.4), size=Pt(9), bold=True, color=color, align=PP_ALIGN.RIGHT)
    # dot
    rect(sl, Inches(2.25), py + Inches(0.42), Inches(0.18), Inches(0.18),
         fill_color=color)
    # connector (except last)
    if i < len(phases) - 1:
        rect(sl, Inches(2.32), py + Inches(0.6), Pt(2), Inches(1.4) + gap_p - Inches(0.18),
             fill_color=RGBColor(0x1f,0x3b,0x38))
    # content card
    rect(sl, Inches(2.6), py, Inches(10.0), phase_h,
         fill_color=CARD_BG, line_color=CARD_BORDER, line_width=Pt(0.5))
    txbox(sl, ptitle, Inches(2.8), py + Inches(0.12),
          Inches(9.6), Inches(0.38), size=Pt(14), bold=True)
    txbox(sl, items, Inches(2.8), py + Inches(0.52),
          Inches(9.6), Inches(0.75), size=Pt(11), color=GRAY_DIM, wrap=True)


# ════════════════════════════════════════════════════════════════════════════
#  SLIDE 14 — CLOSING
# ════════════════════════════════════════════════════════════════════════════
sl = add_slide()
bg(sl, NEAR_BLACK)
rect(sl, Inches(0), Inches(0), W, H, fill_color=RGBColor(0x05,0x12,0x10))

txbox(sl,
      "Every home visit\nis a second chance.",
      Inches(1.0), Inches(1.2), Inches(11.3), Inches(2.5),
      size=Pt(60), bold=True, color=WHITE, align=PP_ALIGN.CENTER)

txbox(sl,
      "AAMA सखी gives Nepal's community health volunteers the tools to catch what\n"
      "the system has always missed — safely, privately, and without internet.",
      Inches(2.0), Inches(3.85), Inches(9.3), Inches(1.2),
      size=Pt(15), color=GRAY_DIM, align=PP_ALIGN.CENTER, wrap=True)

# accent line
rect(sl, Inches(5.5), Inches(5.2), Inches(2.3), Pt(2), fill_color=TEAL_LIGHT)

txbox(sl, "React Native  ·  Expo  ·  SQLite  ·  EPDS  ·  PHQ-A  ·  Offline-First  ·  Bilingual",
      Inches(1.0), Inches(5.5), Inches(11.3), Inches(0.5),
      size=Pt(11), color=GRAY_DIM, align=PP_ALIGN.CENTER)

txbox(sl, "PROBLEM  ·  SOLUTION  ·  IMPACT  ·  TECH  ·  CHALLENGES",
      Inches(1.0), Inches(6.3), Inches(11.3), Inches(0.4),
      size=Pt(9), bold=True, color=RGBColor(0x1f,0x3b,0x38), align=PP_ALIGN.CENTER)


# ── Save ────────────────────────────────────────────────────────────────────
out = "/Users/sujalneupane/Documents/Code/AAMA/AAMA_Pitch_Deck.pptx"
prs.save(out)
print(f"Saved: {out}")
